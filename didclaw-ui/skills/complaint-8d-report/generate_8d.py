#!/usr/bin/env python3
"""
8D Report generator — produces .pptx and/or .pdf from a JSON data file.

Usage:
    python generate_8d.py <input.json> [--format pptx|pdf|both] [--out <dir>]

    --format   Output format: pptx (default), pdf, or both
    --out      Output directory (default: same directory as input file)

The input JSON schema is documented in SKILL.md (section "Generating files").

Dependencies (install once):
    pip install python-pptx fpdf2
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from datetime import date
from pathlib import Path

PLACEHOLDER = "[To be filled]"


def _v(data: dict, key: str, fallback: str = PLACEHOLDER) -> str:
    val = data.get(key)
    return str(val).strip() if val and str(val).strip() else fallback


# ─────────────────────────── PPTX helpers ───────────────────────────

def _title_box(slide, text: str):
    from pptx.util import Inches, Pt
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.65))
    tf = box.text_frame
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(22)
    run.font.bold = True


def _body_box(slide, lines: list[str], top_inches: float = 0.9):
    from pptx.util import Inches, Pt
    box = slide.shapes.add_textbox(Inches(0.4), Inches(top_inches), Inches(9.2), Inches(6.4 - top_inches))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            tf.text = line
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].font.size = Pt(11)
            first = False
        else:
            p = tf.add_paragraph()
            p.text = line
            if p.runs:
                p.runs[0].font.size = Pt(11)


def _table(slide, headers: list[str], rows: list[list[str]], col_widths: list[float], top: float):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    n_rows = len(rows) + 1
    total_w = sum(col_widths)
    tbl = slide.shapes.add_table(
        n_rows, len(headers),
        Inches(0.4), Inches(top),
        Inches(total_w), Inches(0.38 * n_rows),
    ).table

    for j, (h, w) in enumerate(zip(headers, col_widths)):
        tbl.columns[j].width = Inches(w)
        cell = tbl.cell(0, j)
        cell.text = h
        run = cell.text_frame.paragraphs[0].runs
        if run:
            run[0].font.bold = True
            run[0].font.size = Pt(10)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val) if val else ""
            # Ensure there's at least one run to format
            if not cell.text_frame.paragraphs[0].runs:
                cell.text_frame.paragraphs[0].add_run()
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xED, 0xF2, 0xF8)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


# ─────────────────────────── PPTX slides ────────────────────────────

def _slide_cover(prs, data: dict):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    tf.text = "Customer Complaint 8D Report"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    subtitle_lines = [
        f"Case No.:   {_v(data, 'case_no')}",
        f"Product:     {_v(data, 'product')}",
        f"Customer:  {_v(data, 'customer')}",
        f"Date:          {_v(data, 'date', str(date.today()))}",
        f"Owner:        {_v(data, 'owner')}",
    ]
    box2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(8.8), Inches(3.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.text = subtitle_lines[0]
    tf2.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xF0)
    tf2.paragraphs[0].runs[0].font.size = Pt(14)
    for line in subtitle_lines[1:]:
        p = tf2.add_paragraph()
        p.text = line
        if p.runs:
            p.runs[0].font.color.rgb = RGBColor(0xCC, 0xDD, 0xF0)
            p.runs[0].font.size = Pt(14)


def _slide_basic_info(prs, data: dict):
    slide = _blank_slide(prs)
    _title_box(slide, "Basic Information")
    rows = [
        ["Complaint / Case No.", _v(data, "case_no")],
        ["Product name / Model", _v(data, "product")],
        ["Batch / Lot or production date", _v(data, "batch")],
        ["Customer name", _v(data, "customer")],
        ["Complaint date", _v(data, "complaint_date")],
        ["Receipt date", _v(data, "receipt_date", PLACEHOLDER)],
        ["Defect summary", _v(data, "defect_summary")],
        ["Report date", _v(data, "date", str(date.today()))],
        ["Owner / Team leader", _v(data, "owner")],
    ]
    for opt_key, opt_label in [
        ("rma_no", "RMA / Customer complaint No."),
        ("response_due", "Supplier response due date"),
        ("cost_impact", "Cost of quality / impact"),
    ]:
        if data.get(opt_key):
            rows.append([opt_label, str(data[opt_key])])
    _table(slide, ["Item", "Content"], rows, [3.2, 6.0], top=0.9)


def _slide_d1(prs, data: dict):
    slide = _blank_slide(prs)
    _title_box(slide, "D1 – Form the Team")
    team = data.get("d1_team") or [
        {"role": "Leader", "name": PLACEHOLDER, "dept": PLACEHOLDER, "responsibility": "Overall coordination"},
    ]
    rows = [[
        str(m.get("role", "")), str(m.get("name", "")),
        str(m.get("dept", "")), str(m.get("responsibility", "")),
        str(m.get("contact", "")),
    ] for m in team]
    _table(slide, ["Role", "Name", "Department", "Responsibility", "Contact"],
           rows, [1.5, 1.8, 1.8, 2.8, 1.3], top=0.9)


def _slide_d2(prs, data: dict):
    slide = _blank_slide(prs)
    _title_box(slide, "D2 – Problem Description")
    w = data.get("d2_5w2h") or {}
    lines = [
        f"What:          {w.get('what') or PLACEHOLDER}",
        f"Where:         {w.get('where') or PLACEHOLDER}",
        f"When:          {w.get('when') or PLACEHOLDER}",
        f"Who:           {w.get('who') or PLACEHOLDER}",
        f"Why:           {w.get('why') or PLACEHOLDER}",
        f"How many:  {w.get('how_many') or PLACEHOLDER}",
        f"How detected: {w.get('how') or PLACEHOLDER}",
        "",
        f"Problem statement:  {_v(data, 'd2_problem_statement')}",
        "",
        "[Attach: defect photo(s), limit sample or spec ref]",
    ]
    _body_box(slide, lines)


def _slide_d3(prs, data: dict):
    slide = _blank_slide(prs)
    _title_box(slide, "D3 – Interim Containment")
    actions = data.get("d3_actions") or [
        {"no": 1, "action": PLACEHOLDER, "owner": "", "due_date": "", "verification": ""},
    ]
    rows = [[
        str(a.get("no", i + 1)), str(a.get("action", "")),
        str(a.get("owner", "")), str(a.get("due_date", "")),
        str(a.get("verification", "")),
    ] for i, a in enumerate(actions)]
    _table(slide, ["No.", "Action", "Owner", "Due date", "Verification"],
           rows, [0.5, 4.0, 1.5, 1.5, 1.7], top=0.9)


def _slide_d4(prs, data: dict):
    slide = _blank_slide(prs)
    _title_box(slide, "D4 – Root Cause")
    five_why = data.get("d4_5why") or []
    lines: list[str] = []
    if five_why:
        lines.append("5-Why Analysis:")
        for i, why in enumerate(five_why, 1):
            lines.append(f"  Why {i}: {why}")
        lines.append("")
    lines += [
        f"Direct cause:  {_v(data, 'd4_direct_cause')}",
        f"Root cause:    {_v(data, 'd4_root_cause')}",
        f"Verification:  {_v(data, 'd4_verification')}",
        "",
        "[Attach: 5-Why tree / fishbone diagram / Pareto chart]",
        "[Attach: reproduction test photo or video (if available)]",
    ]
    _body_box(slide, lines)


def _slide_d5(prs, data: dict):
    slide = _blank_slide(prs)
    _title_box(slide, "D5 – Permanent Corrective Actions")
    actions = data.get("d5_actions") or [
        {"id": "PC1", "action": PLACEHOLDER, "owner": "", "planned": "", "result": ""},
    ]
    rows = [[
        str(a.get("id", f"PC{i+1}")), str(a.get("action", "")),
        str(a.get("owner", "")), str(a.get("planned", "")),
        str(a.get("result", "")),
    ] for i, a in enumerate(actions)]
    _table(slide, ["ID", "Action", "Owner", "Planned completion", "Verification result"],
           rows, [0.6, 3.8, 1.5, 1.8, 1.5], top=0.9)


def _slide_d6(prs, data: dict):
    slide = _blank_slide(prs)
    _title_box(slide, "D6 – Implement & Validate")
    lines = [
        f"Updated documents:  {_v(data, 'd6_documents')}",
        f"Implementation scope:  {_v(data, 'd6_scope')}",
        f"Effectiveness evidence:  {_v(data, 'd6_effectiveness')}",
        "",
        "[Attach: revised SOP / spec / FMEA — cover or key page with rev and date]",
        "[Attach: approval or change record]",
        "[Attach: Cpk or trend chart (optional)]",
    ]
    _body_box(slide, lines)


def _slide_d7(prs, data: dict):
    slide = _blank_slide(prs)
    _title_box(slide, "D7 – Prevent Recurrence")
    prevention = data.get("d7_prevention") or [
        {"category": "Process / Standard", "content": PLACEHOLDER, "owner": ""},
        {"category": "Horizontal deployment", "content": PLACEHOLDER, "owner": ""},
        {"category": "Training / Sharing", "content": PLACEHOLDER, "owner": ""},
    ]
    rows = [[
        str(p.get("category", "")), str(p.get("content", "")), str(p.get("owner", "")),
    ] for p in prevention]
    _table(slide, ["Category", "Content", "Owner"], rows, [2.5, 5.5, 1.2], top=0.9)


def _slide_d8(prs, data: dict):
    slide = _blank_slide(prs)
    _title_box(slide, "D8 – Congratulate & Close")
    lines = [
        f"Closure summary:  {_v(data, 'd8_closure')}",
        f"Closure date:        {_v(data, 'd8_closure_date')}",
        f"Distribution:          {_v(data, 'd8_distribution')}",
        "",
        "[Attach: customer closure confirmation / sign-off]",
        "[Attach: distribution list or acknowledgment]",
    ]
    _body_box(slide, lines)


def generate_pptx(data: dict, output_path: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs, data)
    _slide_basic_info(prs, data)
    _slide_d1(prs, data)
    _slide_d2(prs, data)
    _slide_d3(prs, data)
    _slide_d4(prs, data)
    _slide_d5(prs, data)
    _slide_d6(prs, data)
    _slide_d7(prs, data)
    _slide_d8(prs, data)

    prs.save(output_path)
    return output_path


# ─────────────────────────── PDF helpers ────────────────────────────

def _detect_cjk_font() -> str | None:
    """Return path to a CJK TTF font if one exists on this machine."""
    candidates = [
        # Windows
        r"C:\Windows\Fonts\msyh.ttf",       # Microsoft YaHei
        r"C:\Windows\Fonts\simhei.ttf",     # SimHei
        r"C:\Windows\Fonts\simsun.ttc",
        # macOS
        "/System/Library/Fonts/PingFang.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        # Linux
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJKsc-Regular.otf",
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    return None


def generate_pdf(data: dict, output_path: str) -> str:
    from fpdf import FPDF

    cjk_font = _detect_cjk_font()

    class PDF(FPDF):
        FONT_REGULAR = "CJK" if cjk_font else "Helvetica"
        FONT_BOLD = "CJK" if cjk_font else "Helvetica"

        def header(self):
            self.set_font(self.FONT_BOLD, size=9)
            case = data.get("case_no", "")
            self.cell(0, 6, f"8D Report  {('— ' + case) if case else ''}", ln=True)
            self.ln(1)

        def footer(self):
            self.set_y(-12)
            self.set_font(self.FONT_REGULAR, size=8)
            self.set_text_color(150)
            self.cell(0, 5, f"Page {self.page_no()}", align="C")
            self.set_text_color(0)

    pdf = PDF()
    pdf.set_auto_page_break(auto=True, margin=16)
    pdf.set_margins(14, 14, 14)

    if cjk_font:
        # fpdf2 supports uni=True for TTF/OTF fonts — handles CJK characters
        pdf.add_font("CJK", style="", fname=cjk_font, uni=True)

    FONT = "CJK" if cjk_font else "Helvetica"

    def section_title(title: str):
        pdf.set_font(FONT, size=14)
        pdf.set_fill_color(31, 73, 125)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(0, 9, f"  {title}", ln=True, fill=True)
        pdf.set_text_color(0, 0, 0)
        pdf.ln(2)

    def kv_row(label: str, value: str):
        pdf.set_font(FONT, size=9)
        pdf.set_fill_color(237, 242, 248)
        pdf.cell(60, 6, label, border=1, fill=True)
        pdf.set_fill_color(255, 255, 255)
        pdf.multi_cell(0, 6, value or PLACEHOLDER, border=1)

    def body_text(text: str):
        pdf.set_font(FONT, size=10)
        pdf.multi_cell(0, 6, text)
        pdf.ln(1)

    def table_rows(headers: list[str], rows: list[list[str]], col_widths: list[float]):
        usable = pdf.w - pdf.l_margin - pdf.r_margin
        widths = [w * usable for w in col_widths]
        pdf.set_font(FONT, size=9)
        pdf.set_fill_color(31, 73, 125)
        pdf.set_text_color(255, 255, 255)
        for h, w in zip(headers, widths):
            pdf.cell(w, 7, h, border=1, fill=True)
        pdf.ln()
        pdf.set_text_color(0, 0, 0)
        for i, row in enumerate(rows):
            if i % 2 == 0:
                pdf.set_fill_color(255, 255, 255)
            else:
                pdf.set_fill_color(237, 242, 248)
            for val, w in zip(row, widths):
                pdf.cell(w, 6, str(val or ""), border=1, fill=True)
            pdf.ln()
        pdf.ln(2)

    # ── Cover page ──
    pdf.add_page()
    pdf.set_font(FONT, size=20)
    pdf.set_fill_color(31, 73, 125)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 14, "Customer Complaint 8D Report", ln=True, fill=True, align="C")
    pdf.set_text_color(0, 0, 0)
    pdf.ln(6)
    for label, key in [
        ("Case No.", "case_no"), ("Product", "product"), ("Customer", "customer"),
        ("Date", "date"), ("Owner", "owner"),
    ]:
        kv_row(label, _v(data, key))
    pdf.ln(4)

    # ── Basic information ──
    section_title("Basic Information")
    for label, key in [
        ("Complaint / Case No.", "case_no"), ("Product name / Model", "product"),
        ("Batch / Lot", "batch"), ("Customer name", "customer"),
        ("Complaint date", "complaint_date"), ("Receipt date", "receipt_date"),
        ("Defect summary", "defect_summary"), ("Report date", "date"),
        ("Owner / Team leader", "owner"),
    ]:
        kv_row(label, _v(data, key))
    for opt_key, opt_label in [
        ("rma_no", "RMA / Customer complaint No."),
        ("response_due", "Supplier response due date"),
        ("cost_impact", "Cost of quality / impact"),
    ]:
        if data.get(opt_key):
            kv_row(opt_label, str(data[opt_key]))
    pdf.ln(4)

    # ── D1 ──
    pdf.add_page()
    section_title("D1 – Form the Team")
    team = data.get("d1_team") or [
        {"role": "Leader", "name": PLACEHOLDER, "dept": PLACEHOLDER,
         "responsibility": "Overall coordination", "contact": ""},
    ]
    rows = [[
        str(m.get("role", "")), str(m.get("name", "")),
        str(m.get("dept", "")), str(m.get("responsibility", "")),
        str(m.get("contact", "")),
    ] for m in team]
    table_rows(["Role", "Name", "Dept.", "Responsibility", "Contact"],
               rows, [0.14, 0.16, 0.16, 0.34, 0.20])

    # ── D2 ──
    section_title("D2 – Problem Description")
    w = data.get("d2_5w2h") or {}
    for label, key in [
        ("What", "what"), ("Where", "where"), ("When", "when"),
        ("Who", "who"), ("Why", "why"), ("How many", "how_many"), ("How detected", "how"),
    ]:
        kv_row(label, w.get(key) or PLACEHOLDER)
    pdf.ln(2)
    kv_row("Problem statement", _v(data, "d2_problem_statement"))
    body_text("[Attach: defect photo(s), limit sample or spec ref]")
    pdf.ln(2)

    # ── D3 ──
    pdf.add_page()
    section_title("D3 – Interim Containment")
    actions = data.get("d3_actions") or [
        {"no": 1, "action": PLACEHOLDER, "owner": "", "due_date": "", "verification": ""},
    ]
    rows = [[
        str(a.get("no", i + 1)), str(a.get("action", "")),
        str(a.get("owner", "")), str(a.get("due_date", "")),
        str(a.get("verification", "")),
    ] for i, a in enumerate(actions)]
    table_rows(["No.", "Action", "Owner", "Due date", "Verification"],
               rows, [0.05, 0.42, 0.18, 0.17, 0.18])

    # ── D4 ──
    section_title("D4 – Root Cause")
    five_why = data.get("d4_5why") or []
    if five_why:
        pdf.set_font(FONT, size=10)
        pdf.cell(0, 6, "5-Why Analysis:", ln=True)
        for i, why in enumerate(five_why, 1):
            pdf.cell(8)
            pdf.multi_cell(0, 6, f"Why {i}: {why}")
        pdf.ln(2)
    for label, key in [
        ("Direct cause", "d4_direct_cause"),
        ("Root cause", "d4_root_cause"),
        ("Verification", "d4_verification"),
    ]:
        kv_row(label, _v(data, key))
    body_text("[Attach: 5-Why tree / fishbone diagram / Pareto chart]")
    pdf.ln(2)

    # ── D5 ──
    pdf.add_page()
    section_title("D5 – Permanent Corrective Actions")
    actions = data.get("d5_actions") or [
        {"id": "PC1", "action": PLACEHOLDER, "owner": "", "planned": "", "result": ""},
    ]
    rows = [[
        str(a.get("id", f"PC{i+1}")), str(a.get("action", "")),
        str(a.get("owner", "")), str(a.get("planned", "")),
        str(a.get("result", "")),
    ] for i, a in enumerate(actions)]
    table_rows(["ID", "Action", "Owner", "Planned completion", "Verification result"],
               rows, [0.06, 0.38, 0.18, 0.20, 0.18])

    # ── D6 ──
    section_title("D6 – Implement & Validate")
    for label, key in [
        ("Updated documents", "d6_documents"),
        ("Implementation scope", "d6_scope"),
        ("Effectiveness evidence", "d6_effectiveness"),
    ]:
        kv_row(label, _v(data, key))
    body_text("[Attach: revised SOP / spec — cover or key page with rev and date, approval record]")
    pdf.ln(2)

    # ── D7 ──
    section_title("D7 – Prevent Recurrence")
    prevention = data.get("d7_prevention") or [
        {"category": "Process / Standard", "content": PLACEHOLDER, "owner": ""},
        {"category": "Horizontal deployment", "content": PLACEHOLDER, "owner": ""},
        {"category": "Training / Sharing", "content": PLACEHOLDER, "owner": ""},
    ]
    rows = [[
        str(p.get("category", "")), str(p.get("content", "")), str(p.get("owner", "")),
    ] for p in prevention]
    table_rows(["Category", "Content", "Owner"], rows, [0.22, 0.62, 0.16])

    # ── D8 ──
    section_title("D8 – Congratulate & Close")
    for label, key in [
        ("Closure summary", "d8_closure"),
        ("Closure date", "d8_closure_date"),
        ("Distribution", "d8_distribution"),
    ]:
        kv_row(label, _v(data, key))
    body_text("[Attach: customer closure confirmation / sign-off, distribution list]")

    pdf.output(output_path)
    return output_path


# ─────────────────────────── CLI entry ──────────────────────────────

def _stem(data: dict) -> str:
    case = data.get("case_no", "").strip().replace(" ", "-") or "unknown"
    d = data.get("date", str(date.today()))
    return f"8D-{case}-{d}"


def main():
    parser = argparse.ArgumentParser(description="Generate 8D Report (.pptx / .pdf)")
    parser.add_argument("input", help="JSON input file (use - for stdin)")
    parser.add_argument("--format", choices=["pptx", "pdf", "both"], default="pptx",
                        help="Output format (default: pptx)")
    parser.add_argument("--out", default=None,
                        help="Output directory (default: same directory as input file)")
    args = parser.parse_args()

    if args.input == "-":
        data = json.load(sys.stdin)
        base_dir = Path.cwd()
    else:
        input_path = Path(args.input).resolve()
        with open(input_path, encoding="utf-8") as f:
            data = json.load(f)
        base_dir = input_path.parent

    out_dir = Path(args.out).resolve() if args.out else base_dir
    out_dir.mkdir(parents=True, exist_ok=True)

    stem = _stem(data)
    generated: list[str] = []

    if args.format in ("pptx", "both"):
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
            sys.exit(1)
        path = str(out_dir / f"{stem}.pptx")
        generate_pptx(data, path)
        generated.append(path)
        print(f"PPTX: {path}")
        uri = "file:///" + path.replace("\\", "/")
        print(f"Open: {uri}")

    if args.format in ("pdf", "both"):
        try:
            from fpdf import FPDF  # noqa: F401
        except ImportError:
            print("ERROR: fpdf2 not installed. Run: pip install fpdf2", file=sys.stderr)
            sys.exit(1)
        path = str(out_dir / f"{stem}.pdf")
        generate_pdf(data, path)
        generated.append(path)
        print(f"PDF:  {path}")
        uri = "file:///" + path.replace("\\", "/")
        print(f"Open: {uri}")

    if not generated:
        print("No files generated.", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
